from __future__ import annotations

import argparse
import html
import math
import re
import sys
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from analisar_microdados_avaliarj_excel import (  # noqa: E402
    build_ranked_reports,
    build_skill_summary,
    infer_component,
    load_legend,
    normalize_skill_code,
    prepare_base,
    read_input,
    weighted_average,
)


NAVY = RGBColor(20, 45, 78)
BLUE = RGBColor(35, 92, 158)
TEAL = RGBColor(33, 140, 145)
GOLD = RGBColor(214, 158, 46)
SAND = RGBColor(243, 239, 232)
INK = RGBColor(38, 48, 64)
RED = RGBColor(190, 77, 63)
GREEN = RGBColor(82, 141, 86)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Gera dashboard HTML e apresentacao PowerPoint a partir de arquivos agregados do AvaliaRJ."
    )
    parser.add_argument("--arquivo-escola", required=True)
    parser.add_argument("--arquivo-turma", required=True)
    parser.add_argument("--saida-html", required=True)
    parser.add_argument("--saida-pptx", required=True)
    parser.add_argument(
        "--arquivo-legenda",
        default="export/generated/legenda_habilidades_alfabetizarj_2o_captura_2025.csv",
    )
    return parser.parse_args()


def pt_number(value: float, decimals: int = 0) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "-"
    formatted = f"{value:,.{decimals}f}"
    return formatted.replace(",", "X").replace(".", ",").replace("X", ".")


def pt_percent(value: float, decimals: int = 1) -> str:
    return f"{pt_number(value, decimals)}%"


def learning_band(value: float) -> dict[str, str]:
    if value < 25:
        return {
            "label": "Muito abaixo do esperado",
            "tone": "#c95c54",
            "detail": "resultado crítico para a série",
        }
    if value < 50:
        return {
            "label": "Atenção alta",
            "tone": "#d8894d",
            "detail": "aprendizado ainda insuficiente",
        }
    if value < 70:
        return {
            "label": "Em consolidação",
            "tone": "#d2b24c",
            "detail": "há base, mas ainda sem consistência",
        }
    return {
        "label": "Bom patamar",
        "tone": "#2f855a",
        "detail": "aprendizado em nível mais favorável",
    }


def normalize_school_name(name: str) -> str:
    mapping = {
        "EM JAYME SICILIANO": "E. M. Jayme Siciliano",
        "E M JAYME SICILIANO": "E. M. Jayme Siciliano",
        "E M HILDA BRAGA": "E. M. Hilda Braga",
        "EM HILDA BRAGA": "E. M. Hilda Braga",
        "E M MARIA SEMEDO DE ANDRADE": "E. M. Maria Semedo de Andrade",
        "EM MARIA SEMEDO DE ANDRADE": "E. M. Maria Semedo de Andrade",
        "E M ANTONIO DE FREITAS": "E. M. Antônio de Freitas",
        "EM ANTONIO DE FREITAS": "E. M. Antônio de Freitas",
        "E M ANTÔNIO DE FREITAS": "E. M. Antônio de Freitas",
        "EM ANTÔNIO DE FREITAS": "E. M. Antônio de Freitas",
    }
    key = " ".join(str(name).upper().split())
    return mapping.get(key, str(name).title().replace(" E ", " e "))


MATRIX_CODE_MEANINGS = {
    "D01": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Relacionar elementos sonoros das palavras com sua representação escrita.",
    },
    "D02": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Ler palavras.",
    },
    "D03": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Ler frases.",
    },
    "D04": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Localizar informações explícitas em textos.",
    },
    "D05": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Reconhecer a finalidade de um texto.",
    },
    "D06": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Inferir o assunto de um texto.",
    },
    "D07": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Inferir informações em textos verbais.",
    },
    "D08": {
        "dominio": "Língua Portuguesa",
        "tipo": "Matriz de referência",
        "evidencia": "Inferir informações em textos que articulam linguagem verbal e não verbal.",
    },
    "D09": {
        "dominio": "Língua Portuguesa",
        "tipo": "Item de resposta construída",
        "evidencia": "Escrever palavras.",
    },
    "D10": {
        "dominio": "Língua Portuguesa",
        "tipo": "Item de resposta construída",
        "evidencia": "Escrever textos.",
    },
}


def build_skill_meaning(code: str, legend_item: dict[str, object]) -> str:
    code_public = str(legend_item.get("codigo_referencia_publica", "") or "").strip()
    meaning = MATRIX_CODE_MEANINGS.get(code_public)
    if meaning:
        return str(meaning["evidencia"])

    if code_public or code:
        return "Significado do descritor não localizado na matriz de referência disponível."

    return "Significado do descritor não localizado na matriz de referência disponível."


def build_skill_display_label(code: str, legend_item: dict[str, object]) -> str:
    normalized_code = str(code).strip().upper()
    match = re.fullmatch(r"H(\d+)", normalized_code)
    visible_code = f"H {int(match.group(1)):02d}" if match else str(code)
    code_public = str(legend_item.get("codigo_referencia_publica", "") or "").strip()
    if code_public:
        return f"{visible_code} ({code_public})"
    return visible_code


def build_context(args: argparse.Namespace) -> dict[str, object]:
    school_raw = read_input(args.arquivo_escola)
    turma_raw = read_input(args.arquivo_turma)

    component = infer_component(school_raw)
    legend = load_legend(args.arquivo_legenda, component)

    school_base, skill_columns_school, _ = prepare_base(school_raw, include_turma=False)
    turma_base, skill_columns_turma, _ = prepare_base(turma_raw, include_turma=True)
    skill_columns = [column for column in skill_columns_school if column in skill_columns_turma]
    if not skill_columns:
        raise ValueError("Nenhuma coluna de habilidade comum foi encontrada nos arquivos informados.")

    legend_map = legend.set_index("habilidade").to_dict("index")
    skill_labels = {}
    skill_meanings = {}
    for column in skill_columns:
        skill_code = normalize_skill_code(column)
        legend_item = legend_map.get(skill_code, {})
        skill_labels[column] = build_skill_display_label(skill_code, legend_item)
        skill_meanings[skill_labels[column]] = build_skill_meaning(skill_code, legend_item)

    school_report, turma_report, network_prof = build_ranked_reports(
        school_base,
        turma_base,
        skill_columns,
        skill_labels,
    )
    skill_summary, school_skills, turma_skills = build_skill_summary(
        school_base,
        turma_base,
        skill_columns,
        skill_labels,
        legend,
    )

    for df in (school_report, turma_report, school_skills, turma_skills):
        if "Escola" in df.columns:
            df["Escola"] = df["Escola"].map(normalize_school_name)

    metadata_columns = [
        "Avaliação",
        "Rede",
        "Ano Escolar",
        "Componente Curricular",
        "Município",
        "Regional",
        "Estado",
    ]
    metadata = {column: str(school_raw.iloc[0][column]) for column in metadata_columns if column in school_raw.columns}

    total_previstos = int(school_report["Previstos"].sum())
    total_avaliados = int(school_report["Avaliados"].sum())
    participation = (total_avaliados / total_previstos) * 100 if total_previstos else math.nan
    critical = weighted_average(school_report["Crítico (%)"], school_report["Avaliados"])
    expected = weighted_average(school_report["Esperado+ (%)"], school_report["Avaliados"])
    below = weighted_average(school_report["Abaixo do básico %"], school_report["Avaliados"])
    basic = weighted_average(school_report["Básico %"], school_report["Avaliados"])
    adequate = weighted_average(school_report["Adequado %"], school_report["Avaliados"])
    advanced = weighted_average(school_report["Avançado %"], school_report["Avaliados"])

    best_school = school_report.iloc[0].to_dict()
    worst_school = school_report.iloc[-1].to_dict()
    best_turma = turma_report.iloc[0].to_dict()
    worst_turma = turma_report.iloc[-1].to_dict()
    priority_skills = skill_summary.head(3).to_dict("records")
    strengths = skill_summary.tail(3).sort_values("Media da rede (%)", ascending=False).to_dict("records")

    school_skill_lookup = school_skills.set_index("Escola").to_dict("index")
    school_profiles = []
    for _, row in school_report.iterrows():
        skill_row = school_skill_lookup[str(row["Escola"])]
        skill_pairs = [
            (column, float(skill_row[column]))
            for column in school_skills.columns
            if column not in {"Escola", "Avaliados", "Media habilidades (%)"}
        ]
        skill_pairs_sorted = sorted(skill_pairs, key=lambda item: item[1])
        adequate = float(row["Esperado+ (%)"])
        school_profiles.append(
            {
                "school": str(row["Escola"]),
                "ranking": int(row["Ranking"]),
                "prof": float(row["Proficiência Média"]),
                "participation": float(row["Participacao %"]),
                "adequate": adequate,
                "critical": float(row["Crítico (%)"]),
                "delta": float(row["Proficiência vs rede"]),
                "band": learning_band(adequate),
                "distribution": [
                    ("Abaixo do básico", float(row["Abaixo do básico %"])),
                    ("Básico", float(row["Básico %"])),
                    ("Adequado", float(row["Adequado %"])),
                    ("Avançado", float(row["Avançado %"])),
                ],
                "weak_skills": skill_pairs_sorted[:3],
                "strong_skills": list(reversed(skill_pairs_sorted[-3:])),
            }
        )

    recommendations = [
        {
            "title": "Foco imediato de recomposição",
            "body": (
                "Organizar intervenções de leitura e compreensão com foco nas habilidades "
                + ", ".join(item["Habilidade"] for item in priority_skills)
                + "."
            ),
        },
        {
            "title": "Difusão de prática de referência",
            "body": (
                f"Sistematizar estratégias da escola {best_school['Escola']} e, em especial, "
                f"da turma {best_turma['Turma']} para compartilhar rotinas de sucesso na rede."
            ),
        },
        {
            "title": "Monitoramento intensivo",
            "body": (
                f"Acompanhar quinzenalmente a turma {worst_turma['Turma']} da escola "
                f"{worst_turma['Escola']}, que concentrou a menor proficiência do recorte."
            ),
        },
        {
            "title": "Sustentar a presença",
            "body": (
                f"Manter a participação acima de 90% e buscar fechamento do gap remanescente "
                f"de {pt_percent(100 - participation)} entre previstos e avaliados."
            ),
        },
    ]

    return {
        "metadata": metadata,
        "school_report": school_report,
        "turma_report": turma_report,
        "skill_summary": skill_summary,
        "school_skills": school_skills,
        "turma_skills": turma_skills,
        "network_prof": network_prof,
        "total_previstos": total_previstos,
        "total_avaliados": total_avaliados,
        "participation": participation,
        "critical": critical,
        "expected": expected,
        "levels": {
            "Abaixo do básico": below,
            "Básico": basic,
            "Adequado": adequate,
            "Avançado": advanced,
        },
        "best_school": best_school,
        "worst_school": worst_school,
        "best_turma": best_turma,
        "worst_turma": worst_turma,
        "priority_skills": priority_skills,
        "strengths": strengths,
        "school_profiles": school_profiles,
        "skill_meanings": skill_meanings,
        "analysis_scope": "Língua Portuguesa · bloco principal H01-H08",
        "analysis_scope_note": (
            "Os dois arquivos atuais trazem o bloco principal de Língua Portuguesa "
            "e não exibem os itens de resposta construída D09 e D10 neste recorte. "
            "Ao passar o mouse sobre as habilidades, o tooltip exibe apenas o significado do descritor."
        ),
        "recommendations": recommendations,
        "generated_at": datetime.now(),
        "source_escola": str(Path(args.arquivo_escola).resolve()),
        "source_turma": str(Path(args.arquivo_turma).resolve()),
    }


def heat_color(value: float) -> str:
    if value >= 85:
        return "#2f855a"
    if value >= 70:
        return "#4c9f70"
    if value >= 60:
        return "#d2b24c"
    if value >= 50:
        return "#d8894d"
    return "#c95c54"


def neutral_bar(value: float, color: str) -> str:
    return (
        '<div class="bar-track"><span class="bar-fill" '
        f'style="width:{max(0.0, min(value, 100.0)):.1f}%; background:{color};"></span></div>'
    )


def adequacy_scale_html(value: float) -> str:
    marker_left = max(0.0, min(value, 100.0))
    return f"""
    <div class="adequacy-scale">
      <div class="adequacy-scale__segments">
        <span class="seg seg-1"></span>
        <span class="seg seg-2"></span>
        <span class="seg seg-3"></span>
        <span class="seg seg-4"></span>
      </div>
      <div class="adequacy-scale__marker" style="left:{marker_left:.1f}%"></div>
      <div class="adequacy-scale__ticks">
        <span>0</span><span>25</span><span>50</span><span>70</span><span>100</span>
      </div>
    </div>
    """


def stacked_levels_html(levels: list[tuple[str, float]]) -> str:
    palette = {
        "Abaixo do básico": "#c95c54",
        "Básico": "#d8894d",
        "Adequado": "#d2b24c",
        "Avançado": "#2f855a",
    }
    segments = []
    legends = []
    for label, value in levels:
        segments.append(
            f'<span class="level-stack__segment" style="width:{max(0.0, min(value, 100.0)):.1f}%; background:{palette[label]};"></span>'
        )
        legends.append(
            f'<div class="level-legend__item"><i style="background:{palette[label]};"></i><span>{html.escape(label)}</span><strong>{pt_percent(value)}</strong></div>'
        )
    return f"""
    <div class="level-stack">{''.join(segments)}</div>
    <div class="level-legend">{''.join(legends)}</div>
    """


def render_skill_label_html(label: str) -> str:
    label_text = str(label).strip()
    match = re.fullmatch(r"(H \d{2}) \((D\d{2})\)", label_text)
    if not match:
        return f"<span>{html.escape(label_text)}</span>"
    return (
        '<span class="skill-chip__label">'
        f'<span class="skill-chip__main">{html.escape(match.group(1))}</span>'
        f'<span class="skill-chip__sub">({html.escape(match.group(2))})</span>'
        "</span>"
    )


def tooltip_skill_html(label: str, value: float | None, meaning: str, tone: str) -> str:
    score = f"<strong>{pt_percent(value)}</strong>" if value is not None else ""
    label_html = render_skill_label_html(label)
    return (
        f'<button type="button" class="skill-chip {tone}" aria-label="{html.escape(str(label), quote=True)}" data-tooltip="{html.escape(meaning, quote=True)}">'
        f"{label_html}{score}</button>"
    )


def heat_text_color(value: float) -> str:
    if 50 <= value < 70:
        return "#263040"
    return "#ffffff"


def heat_tile_html(label: str, value: float, meaning: str) -> str:
    return (
        '<button type="button" class="heat-tile" '
        f'aria-label="{html.escape(str(label), quote=True)}" '
        f'data-tooltip="{html.escape(meaning, quote=True)}" '
        f'style="background:{heat_color(value)}; color:{heat_text_color(value)};">'
        f'{render_skill_label_html(label)}'
        f"<strong>{pt_percent(value)}</strong>"
        "</button>"
    )


def heatmap_card_html(
    title: str,
    subtitle: str,
    meta_items: list[tuple[str, str]],
    skill_items: list[tuple[str, float, str]],
) -> str:
    meta_html = "".join(
        f'<div class="heatmap-card__meta-item"><span>{html.escape(label)}</span><strong>{html.escape(value)}</strong></div>'
        for label, value in meta_items
    )
    tiles_html = "".join(
        heat_tile_html(label, value, meaning)
        for label, value, meaning in skill_items
    )
    return f"""
    <article class="heatmap-card">
      <div class="heatmap-card__head">
        <div>
          <h4>{html.escape(title)}</h4>
          <p>{html.escape(subtitle)}</p>
        </div>
        <div class="heatmap-card__meta">{meta_html}</div>
      </div>
      <div class="heatmap-card__grid">{tiles_html}</div>
    </article>
    """


def school_turma_card_html(
    row: pd.Series,
    local_rank: int,
    skill_meanings: dict[str, str],
    highlight: str | None = None,
) -> str:
    fragile_label = str(row["Habilidade mais frágil"])
    fragile_meaning = skill_meanings.get(fragile_label, fragile_label)
    band = learning_band(float(row["Esperado+ (%)"]))
    status_style = f"background:{band['tone']}22; color:{band['tone']}; border-color:{band['tone']}33;"
    highlight_html = ""
    if highlight == "best":
        highlight_html = '<span class="turma-pill turma-pill--best">Melhor da escola</span>'
    elif highlight == "worst":
        highlight_html = '<span class="turma-pill turma-pill--worst">Maior atenção</span>'
    return f"""
    <article class="turma-card" style="--turma-accent:{band['tone']};">
      <div class="turma-card__badges">
        <span class="turma-pill turma-pill--status" style="{status_style}">{html.escape(band["label"])}</span>
        {highlight_html}
      </div>
      <div class="turma-card__head">
        <div class="rank-chip">{local_rank}</div>
        <div>
          <h4>{html.escape(str(row["Turma"]))}</h4>
          <p>{html.escape(band["detail"])}</p>
        </div>
      </div>
      <div class="turma-card__metrics">
        <div><span>Proficiência</span><strong>{pt_number(float(row["Proficiência Média"]))}</strong></div>
        <div><span>Participação</span><strong>{pt_percent(float(row["Participacao %"]))}</strong></div>
        <div><span>Esperado+</span><strong>{pt_percent(float(row["Esperado+ (%)"]))}</strong></div>
      </div>
      <div class="turma-card__focus">
        <span class="mini-label">Habilidade mais frágil</span>
        {tooltip_skill_html(fragile_label, None, fragile_meaning, "tag")}
      </div>
    </article>
    """


def render_dashboard_html(context: dict[str, object], output_path: Path) -> Path:
    metadata = context["metadata"]
    school_report = context["school_report"]
    turma_report = context["turma_report"]
    skill_summary = context["skill_summary"]
    school_skills = context["school_skills"]
    turma_skills = context["turma_skills"]
    school_profiles = context["school_profiles"]
    skill_meanings = context["skill_meanings"]

    school_skill_columns = [str(col) for col in school_skills.columns if str(col) in skill_meanings]
    turma_skill_columns = [str(col) for col in turma_skills.columns if str(col) in skill_meanings]
    network_skill_map = skill_summary.set_index("Habilidade")["Media da rede (%)"].to_dict()
    network_heat_card = heatmap_card_html(
        f"Rede municipal de {metadata['Município']}",
        "Domínio médio ponderado por habilidade",
        [
            ("Avaliados", pt_number(float(context["total_avaliados"]))),
            ("Proficiência", pt_number(float(context["network_prof"]), 2)),
            ("Esperado+", pt_percent(float(context["expected"]))),
        ],
        [
            (
                label,
                float(network_skill_map[label]),
                skill_meanings.get(label, label),
            )
            for label in school_skill_columns
            if label in network_skill_map
        ],
    )

    school_tabs = []
    school_panels = []
    for index, profile in enumerate(school_profiles):
        panel_id = f"school-panel-{profile['ranking']}"
        tab_id = f"school-tab-{profile['ranking']}"
        active_class = " active" if index == 0 else ""
        selected = "true" if index == 0 else "false"
        school_name = str(profile["school"])
        school_turma_report = turma_report.loc[turma_report["Escola"] == school_name].copy()
        school_turma_skills = turma_skills.loc[turma_skills["Escola"] == school_name].copy()
        weak_skills = "".join(
            tooltip_skill_html(label, value, skill_meanings.get(label, label), "weak")
            for label, value in profile["weak_skills"]
        )
        strong_skills = "".join(
            tooltip_skill_html(label, value, skill_meanings.get(label, label), "strong")
            for label, value in profile["strong_skills"]
        )
        school_turma_report = school_turma_report.sort_values(
            by=["Proficiência Média", "Esperado+ (%)", "Participacao %", "Turma"],
            ascending=[False, False, False, True],
        ).reset_index(drop=True)

        school_turma_cards = []
        last_index = len(school_turma_report) - 1
        for local_rank, (_, turma_row) in enumerate(school_turma_report.iterrows(), start=1):
            highlight = None
            current_index = local_rank - 1
            if last_index > 0:
                if current_index == 0:
                    highlight = "best"
                elif current_index == last_index:
                    highlight = "worst"
            school_turma_cards.append(school_turma_card_html(turma_row, local_rank, skill_meanings, highlight))

        school_turma_heat_cards = []
        for _, turma_skill_row in school_turma_skills.iterrows():
            skill_items = [
                (
                    col,
                    float(turma_skill_row[col]),
                    skill_meanings.get(col, col),
                )
                for col in turma_skill_columns
            ]
            school_turma_heat_cards.append(
                heatmap_card_html(
                    str(turma_skill_row["Turma"]),
                    "Domínio por habilidade da turma",
                    [
                        ("Avaliados", pt_number(float(turma_skill_row["Avaliados"]))),
                        ("Média", pt_percent(float(turma_skill_row["Media habilidades (%)"]))),
                    ],
                    skill_items,
                )
            )
        school_tabs.append(
            f"""
            <button
              type="button"
              class="school-tab{active_class}"
              id="{tab_id}"
              role="tab"
              aria-selected="{selected}"
              aria-controls="{panel_id}"
              data-school-tab="{panel_id}">
              <span class="school-tab__meta">Escola #{profile["ranking"]}</span>
              <strong>{html.escape(profile["school"])}</strong>
              <small>Esperado+ {pt_percent(profile["adequate"])}</small>
            </button>
            """
        )
        school_panels.append(
            f"""
            <div
              class="school-panel{active_class}"
              id="{panel_id}"
              role="tabpanel"
              aria-labelledby="{tab_id}"
              data-school-panel="{panel_id}"{" hidden" if index != 0 else ""}>
              <article class="school-profile">
                <div class="school-profile__head">
                  <div>
                    <p class="mini-label">Escola #{profile["ranking"]}</p>
                    <h3>{html.escape(profile["school"])}</h3>
                  </div>
                  <div class="school-profile__band" style="background:{profile["band"]["tone"]};">{html.escape(profile["band"]["label"])}</div>
                </div>
                <div class="school-profile__topline">
                  <div class="adequacy-block">
                    <span>Aprendizado adequado</span>
                    <strong>{pt_percent(profile["adequate"])}</strong>
                    <small>{html.escape(profile["band"]["detail"])}</small>
                  </div>
                  <div class="school-meta">
                    <div><span>Proficiência</span><strong>{pt_number(profile["prof"])}</strong></div>
                    <div><span>Participação</span><strong>{pt_percent(profile["participation"])}</strong></div>
                    <div><span>Crítico</span><strong>{pt_percent(profile["critical"])}</strong></div>
                    <div><span>Vs. rede</span><strong>{pt_number(profile["delta"], 2)}</strong></div>
                  </div>
                </div>
                {adequacy_scale_html(profile["adequate"])}
                <div class="school-profile__distribution">
                  <p class="mini-label">Distribuição do aprendizado</p>
                  {stacked_levels_html(profile["distribution"])}
                </div>
                <div class="school-profile__skills">
                  <div>
                    <p class="mini-label">Habilidades que mais precisam de reforço</p>
                    <div class="skill-chip-row">{weak_skills}</div>
                  </div>
                  <div>
                    <p class="mini-label">Habilidades mais consolidadas</p>
                    <div class="skill-chip-row">{strong_skills}</div>
                  </div>
                </div>
              </article>
              <div class="school-detail-panel">
                <p class="mini-label">Análise das turmas da escola</p>
                <div class="turma-card-grid school-turma-card-grid">
                  {''.join(school_turma_cards)}
                </div>
              </div>
              <div class="school-detail-panel">
                <p class="mini-label">Heatmap das turmas da escola</p>
                <div class="heatmap-grid heatmap-grid--turma">
                  {''.join(school_turma_heat_cards)}
                </div>
              </div>
            </div>
            """
        )

    priority_rows = []
    for item in context["priority_skills"]:
        value = float(item["Media da rede (%)"])
        priority_rows.append(
            f"""
            <div class="priority-row">
              <div>
                <p class="mini-label">Prioridade</p>
                {tooltip_skill_html(str(item["Habilidade"]), None, skill_meanings.get(str(item["Habilidade"]), str(item["Habilidade"])), "priority")}
              </div>
              <div class="priority-score">
                <strong>{pt_percent(value)}</strong>
                {neutral_bar(value, heat_color(value))}
              </div>
            </div>
            """
        )

    strength_rows = []
    for item in context["strengths"]:
        value = float(item["Media da rede (%)"])
        strength_rows.append(
            f"""
            {tooltip_skill_html(str(item["Habilidade"]), value, skill_meanings.get(str(item["Habilidade"]), str(item["Habilidade"])), "strong")}
            """
        )


    level_blocks = []
    level_colors = {
        "Abaixo do básico": "#c95c54",
        "Básico": "#d68d4a",
        "Adequado": "#4b8f8c",
        "Avançado": "#2f6f64",
    }
    for label, value in context["levels"].items():
        level_blocks.append(
            f"""
            <div class="level-item">
              <div class="level-item__top">
                <span>{html.escape(label)}</span>
                <strong>{pt_percent(float(value))}</strong>
              </div>
              {neutral_bar(float(value), level_colors[label])}
            </div>
            """
        )

    recommendations_html = []
    for item in context["recommendations"]:
        recommendations_html.append(
            f"""
            <article class="action-card">
              <h4>{html.escape(item["title"])}</h4>
              <p>{html.escape(item["body"])}</p>
            </article>
            """
        )

    html_text = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>AvaliaRJ 2025 · Mendes</title>
  <style>
    :root {{
      --navy:#142d4e;
      --blue:#2d5d8f;
      --teal:#2b7a78;
      --gold:#d6aa3d;
      --sand:#f5f1ea;
      --paper:#fffdf9;
      --ink:#263040;
      --muted:#617084;
      --line:rgba(20,45,78,.12);
      --shadow:0 20px 45px rgba(18,33,56,.12);
      --radius:28px;
    }}
    * {{ box-sizing:border-box; }}
    body {{
      margin:0;
      font-family:"Aptos","Segoe UI","Trebuchet MS",sans-serif;
      color:var(--ink);
      background:
        radial-gradient(circle at top left, rgba(214,170,61,.22), transparent 32%),
        radial-gradient(circle at top right, rgba(43,122,120,.18), transparent 28%),
        linear-gradient(180deg, #f4efe6 0%, #f8f6f1 48%, #ffffff 100%);
    }}
    .page {{
      width:min(1280px, calc(100vw - 48px));
      margin:24px auto 48px;
    }}
    .hero {{
      position:relative;
      overflow:hidden;
      display:grid;
      grid-template-columns:1.25fr .95fr;
      gap:24px;
      padding:34px;
      border-radius:34px;
      background:linear-gradient(135deg, rgba(20,45,78,.98), rgba(35,92,158,.94));
      box-shadow:var(--shadow);
      color:white;
    }}
    .hero::after {{
      content:"";
      position:absolute;
      inset:auto -80px -120px auto;
      width:320px;
      height:320px;
      border-radius:50%;
      background:rgba(214,170,61,.18);
    }}
    .eyebrow {{
      text-transform:uppercase;
      letter-spacing:.18em;
      font-size:12px;
      opacity:.78;
      margin:0 0 14px;
    }}
    h1,h2,h3,h4,p {{ margin:0; }}
    .hero h1 {{
      font-family:"Aptos Display","Bahnschrift","Segoe UI",sans-serif;
      font-size:46px;
      line-height:1.05;
      max-width:700px;
      margin-bottom:16px;
    }}
    .hero p.lead {{
      max-width:720px;
      font-size:17px;
      line-height:1.6;
      color:rgba(255,255,255,.84);
    }}
    .hero-tags {{
      display:flex;
      flex-wrap:wrap;
      gap:10px;
      margin-top:18px;
    }}
    .hero-tags span {{
      padding:10px 14px;
      border-radius:999px;
      border:1px solid rgba(255,255,255,.18);
      background:rgba(255,255,255,.08);
      font-size:13px;
    }}
    .hero-panel {{
      display:grid;
      grid-template-columns:repeat(2,1fr);
      gap:14px;
      align-content:start;
    }}
    .kpi {{
      padding:18px;
      border-radius:22px;
      background:rgba(255,255,255,.11);
      border:1px solid rgba(255,255,255,.14);
      backdrop-filter:blur(8px);
    }}
    .kpi span {{
      display:block;
      font-size:12px;
      text-transform:uppercase;
      letter-spacing:.14em;
      opacity:.72;
      margin-bottom:10px;
    }}
    .kpi strong {{
      display:block;
      font-size:28px;
      line-height:1;
      margin-bottom:6px;
    }}
    .kpi small {{
      color:rgba(255,255,255,.78);
      font-size:13px;
    }}
    .nav {{
      display:flex;
      gap:12px;
      flex-wrap:wrap;
      margin:18px 0 24px;
      padding:10px 0;
    }}
    .nav a {{
      text-decoration:none;
      color:var(--navy);
      background:rgba(255,255,255,.72);
      border:1px solid var(--line);
      padding:10px 14px;
      border-radius:999px;
      font-weight:600;
    }}
    .section {{
      margin-top:20px;
      padding:26px;
      border-radius:30px;
      background:rgba(255,255,255,.82);
      border:1px solid var(--line);
      box-shadow:var(--shadow);
    }}
    .section-head {{
      display:flex;
      justify-content:space-between;
      gap:16px;
      align-items:end;
      margin-bottom:20px;
    }}
    .section-head h2 {{
      font-size:28px;
      color:var(--navy);
      margin-bottom:4px;
    }}
    .section-head p {{
      color:var(--muted);
      line-height:1.5;
      max-width:700px;
    }}
    .overview-grid, .school-grid, .actions-grid {{
      display:grid;
      gap:18px;
    }}
    .overview-grid {{
      grid-template-columns:1.1fr .9fr;
    }}
    .school-tabs {{
      display:grid;
      gap:18px;
    }}
    .school-tablist {{
      display:grid;
      grid-template-columns:repeat(auto-fit, minmax(220px, 1fr));
      gap:12px;
    }}
    .school-tab {{
      padding:16px 18px;
      border-radius:24px;
      border:1px solid rgba(20,45,78,.1);
      background:linear-gradient(180deg, #fffefb 0%, #f3eee4 100%);
      text-align:left;
      cursor:pointer;
      font-family:inherit;
      transition:transform .18s ease, box-shadow .18s ease, background .18s ease;
    }}
    .school-tab:hover {{
      transform:translateY(-1px);
      box-shadow:0 12px 24px rgba(18,33,56,.08);
    }}
    .school-tab:focus-visible {{
      outline:2px solid var(--navy);
      outline-offset:2px;
    }}
    .school-tab.active {{
      background:linear-gradient(135deg, rgba(20,45,78,.98), rgba(35,92,158,.94));
      border-color:rgba(20,45,78,.18);
      box-shadow:0 18px 30px rgba(18,33,56,.14);
    }}
    .school-tab__meta {{
      display:block;
      font-size:11px;
      letter-spacing:.14em;
      text-transform:uppercase;
      color:var(--muted);
      margin-bottom:8px;
    }}
    .school-tab strong {{
      display:block;
      color:var(--navy);
      font-size:17px;
      line-height:1.25;
      margin-bottom:6px;
    }}
    .school-tab small {{
      color:var(--muted);
      font-size:12px;
    }}
    .school-tab.active .school-tab__meta,
    .school-tab.active strong,
    .school-tab.active small {{
      color:white;
    }}
    .school-panels {{
      position:relative;
    }}
    .school-panel {{
      display:grid;
      gap:16px;
    }}
    .school-panel[hidden] {{
      display:none;
    }}
    .school-panel.active {{
      display:grid;
    }}
    .levels-panel, .story-panel, .priority-panel, .heat-panel {{
      border-radius:24px;
      padding:22px;
      background:var(--paper);
      border:1px solid var(--line);
    }}
    .level-item + .level-item {{ margin-top:14px; }}
    .level-item__top, .metric-line {{
      display:flex;
      justify-content:space-between;
      gap:12px;
      align-items:center;
      margin-bottom:8px;
    }}
    .bar-track {{
      width:100%;
      height:12px;
      border-radius:999px;
      overflow:hidden;
      background:#e5e7eb;
    }}
    .bar-fill {{
      display:block;
      height:100%;
      border-radius:999px;
    }}
    .story-list {{
      display:grid;
      gap:16px;
      margin-top:10px;
    }}
    .story-item {{
      padding:16px 18px;
      border-radius:20px;
      background:#f7f3eb;
      border:1px solid rgba(20,45,78,.08);
    }}
    .story-item strong {{
      display:block;
      color:var(--navy);
      margin-bottom:6px;
    }}
    .school-grid {{
      grid-template-columns:1fr;
    }}
    .school-profile {{
      padding:22px;
      border-radius:26px;
      background:linear-gradient(180deg, #fffefb 0%, #f7f1e6 100%);
      border:1px solid rgba(20,45,78,.08);
      box-shadow:0 12px 28px rgba(18,33,56,.06);
    }}
    .school-profile__head {{
      display:flex;
      justify-content:space-between;
      gap:12px;
      align-items:start;
      margin-bottom:14px;
    }}
    .school-profile__head h3 {{
      font-size:20px;
      color:var(--navy);
    }}
    .school-profile__band {{
      color:white;
      font-size:12px;
      font-weight:700;
      padding:10px 14px;
      border-radius:999px;
      white-space:nowrap;
    }}
    .school-profile__topline {{
      display:grid;
      grid-template-columns:.9fr 1.1fr;
      gap:18px;
      margin-bottom:16px;
    }}
    .adequacy-block {{
      padding:18px;
      border-radius:22px;
      background:rgba(255,255,255,.88);
      border:1px solid rgba(20,45,78,.08);
    }}
    .adequacy-block span {{
      display:block;
      font-size:12px;
      text-transform:uppercase;
      letter-spacing:.14em;
      color:var(--muted);
      margin-bottom:10px;
    }}
    .adequacy-block strong {{
      display:block;
      font-size:42px;
      line-height:1;
      color:var(--navy);
      margin-bottom:8px;
    }}
    .adequacy-block small {{
      color:var(--muted);
      font-size:13px;
    }}
    .school-meta {{
      display:grid;
      grid-template-columns:repeat(2,1fr);
      gap:12px;
    }}
    .school-meta div {{
      padding:14px;
      border-radius:18px;
      background:rgba(255,255,255,.88);
      border:1px solid rgba(20,45,78,.08);
    }}
    .school-meta span {{
      display:block;
      color:var(--muted);
      font-size:12px;
      margin-bottom:6px;
    }}
    .school-meta strong {{
      color:var(--navy);
      font-size:18px;
    }}
    .adequacy-scale {{
      position:relative;
      margin-bottom:18px;
      padding-top:8px;
    }}
    .adequacy-scale__segments {{
      display:grid;
      grid-template-columns:25fr 25fr 20fr 30fr;
      gap:4px;
      height:14px;
      border-radius:999px;
      overflow:hidden;
      background:#edf2f7;
    }}
    .seg-1 {{ background:#c95c54; }}
    .seg-2 {{ background:#d8894d; }}
    .seg-3 {{ background:#d2b24c; }}
    .seg-4 {{ background:#2f855a; }}
    .adequacy-scale__marker {{
      position:absolute;
      top:0;
      width:4px;
      height:24px;
      border-radius:999px;
      background:var(--navy);
      transform:translateX(-50%);
      box-shadow:0 0 0 3px rgba(255,255,255,.8);
    }}
    .adequacy-scale__ticks {{
      display:flex;
      justify-content:space-between;
      color:var(--muted);
      font-size:11px;
      margin-top:6px;
    }}
    .school-profile__distribution {{
      padding:16px 18px;
      border-radius:22px;
      background:rgba(255,255,255,.88);
      border:1px solid rgba(20,45,78,.08);
      margin-bottom:16px;
    }}
    .level-stack {{
      display:flex;
      width:100%;
      height:18px;
      border-radius:999px;
      overflow:hidden;
      background:#edf2f7;
      margin-bottom:12px;
    }}
    .level-stack__segment {{
      display:block;
      height:100%;
    }}
    .level-legend {{
      display:grid;
      grid-template-columns:repeat(2,1fr);
      gap:10px 16px;
    }}
    .level-legend__item {{
      display:flex;
      align-items:center;
      gap:8px;
      font-size:12px;
      color:var(--ink);
    }}
    .level-legend__item i {{
      width:10px;
      height:10px;
      border-radius:50%;
      display:inline-block;
      flex:0 0 auto;
    }}
    .level-legend__item strong {{
      margin-left:auto;
      color:var(--navy);
    }}
    .school-profile__skills {{
      display:grid;
      grid-template-columns:1fr 1fr;
      gap:16px;
    }}
    .school-detail-panel {{
      padding:20px;
      border-radius:24px;
      background:var(--paper);
      border:1px solid var(--line);
      box-shadow:0 12px 28px rgba(18,33,56,.05);
    }}
    .skill-chip-row {{
      display:flex;
      flex-wrap:wrap;
      gap:10px;
      margin-top:10px;
    }}
    .skill-chip {{
      display:inline-flex;
      align-items:center;
      gap:8px;
      padding:10px 12px;
      border-radius:999px;
      font-size:12px;
      font-weight:700;
      border:none;
      cursor:pointer;
      position:relative;
      text-align:left;
      font-family:inherit;
    }}
    .skill-chip__label {{
      display:inline-flex;
      align-items:center;
      gap:6px;
      white-space:nowrap;
    }}
    .skill-chip__main {{
      font-weight:800;
      letter-spacing:.02em;
    }}
    .skill-chip__sub {{
      font-size:11px;
      font-weight:700;
      opacity:.72;
    }}
    .skill-chip.header {{
      background:#eef3f8;
      color:var(--navy);
      border:1px solid rgba(20,45,78,.1);
    }}
    .skill-chip.priority {{
      background:#fff6ea;
      color:#8a4a06;
      border:1px solid rgba(216,137,77,.22);
      margin-top:4px;
    }}
    .skill-chip.weak {{
      background:#fff1ef;
      color:#8b2e27;
      border:1px solid rgba(201,92,84,.18);
    }}
    .skill-chip.strong {{
      background:#eef8f1;
      color:#256146;
      border:1px solid rgba(47,133,90,.18);
    }}
    .skill-chip.tag {{
      background:#edf2f7;
      color:var(--navy);
      border:1px solid rgba(20,45,78,.08);
      padding:6px 10px;
      font-size:12px;
      font-weight:600;
    }}
    .skill-chip::after {{
      content:attr(data-tooltip);
      position:absolute;
      left:0;
      top:calc(100% + 10px);
      width:260px;
      padding:12px 14px;
      border-radius:16px;
      background:rgba(20,45,78,.96);
      color:white;
      font-size:12px;
      line-height:1.5;
      box-shadow:0 16px 30px rgba(18,33,56,.22);
      opacity:0;
      transform:translateY(4px);
      pointer-events:none;
      transition:opacity .18s ease, transform .18s ease;
      z-index:20;
      white-space:normal;
    }}
    .skill-chip:hover::after,
    .skill-chip:focus::after,
    .skill-chip:focus-visible::after {{
      opacity:1;
      transform:translateY(0);
    }}
    .skill-chip:focus-visible {{
      outline:2px solid var(--navy);
      outline-offset:2px;
    }}
    .mini-label {{
      font-size:11px;
      letter-spacing:.14em;
      text-transform:uppercase;
      color:var(--muted);
      margin-bottom:6px;
    }}
    .rank-chip {{
      min-width:48px;
      height:48px;
      border-radius:16px;
      display:grid;
      place-items:center;
      background:var(--navy);
      color:white;
      font-weight:700;
    }}
    .table-shell {{
      overflow:auto;
      border-radius:22px;
      border:1px solid rgba(20,45,78,.08);
      background:white;
    }}
    .turma-card-grid {{
      display:grid;
      gap:14px;
    }}
    .school-turma-card-grid {{
      grid-template-columns:repeat(auto-fit, minmax(240px, 1fr));
    }}
    table {{
      width:100%;
      border-collapse:collapse;
      min-width:820px;
    }}
    th, td {{
      padding:14px 12px;
      border-bottom:1px solid rgba(20,45,78,.08);
      text-align:left;
      vertical-align:top;
    }}
    th {{
      position:sticky;
      top:0;
      background:#eef3f8;
      color:var(--navy);
      font-size:13px;
    }}
    .tag {{
      display:inline-flex;
      align-items:center;
      border-radius:999px;
      background:#edf2f7;
      color:var(--navy);
      padding:6px 10px;
      font-size:12px;
      font-weight:600;
    }}
    .priority-grid {{
      display:grid;
      grid-template-columns:1.1fr .9fr;
      gap:18px;
    }}
    .priority-row + .priority-row {{ margin-top:16px; }}
    .priority-row {{
      display:grid;
      grid-template-columns:1.4fr .8fr;
      gap:18px;
      padding:18px;
      border-radius:22px;
      background:#fffefb;
      border:1px solid rgba(20,45,78,.08);
    }}
    .priority-row p {{
      color:var(--muted);
      line-height:1.5;
      margin-top:6px;
    }}
    .priority-score {{
      display:flex;
      flex-direction:column;
      justify-content:center;
      gap:10px;
    }}
    .priority-score strong {{
      font-size:28px;
      color:var(--navy);
    }}
    .turma-card {{
      position:relative;
      overflow:hidden;
      padding:18px;
      border-radius:22px;
      background:#fffefb;
      border:1px solid rgba(20,45,78,.08);
      box-shadow:0 12px 28px rgba(18,33,56,.06);
    }}
    .turma-card::before {{
      content:"";
      position:absolute;
      inset:0 auto 0 0;
      width:6px;
      background:var(--turma-accent, var(--blue));
    }}
    .turma-card__badges {{
      display:flex;
      flex-wrap:wrap;
      gap:8px;
      margin-bottom:14px;
      padding-left:8px;
    }}
    .turma-pill {{
      display:inline-flex;
      align-items:center;
      justify-content:center;
      min-height:28px;
      padding:6px 10px;
      border-radius:999px;
      border:1px solid transparent;
      font-size:11px;
      font-weight:800;
      letter-spacing:.04em;
      text-transform:uppercase;
    }}
    .turma-pill--status {{
      backdrop-filter:blur(8px);
    }}
    .turma-pill--best {{
      background:rgba(43,122,120,.12);
      border-color:rgba(43,122,120,.2);
      color:#1f6d6a;
    }}
    .turma-pill--worst {{
      background:rgba(198,91,82,.12);
      border-color:rgba(198,91,82,.2);
      color:#a94442;
    }}
    .turma-card__head {{
      display:grid;
      grid-template-columns:auto 1fr;
      gap:14px;
      align-items:center;
      margin-bottom:14px;
    }}
    .turma-card__head h4 {{
      color:var(--navy);
      font-size:18px;
      margin-bottom:4px;
    }}
    .turma-card__head p {{
      color:var(--muted);
      line-height:1.4;
    }}
    .turma-card__metrics {{
      display:grid;
      grid-template-columns:repeat(3, 1fr);
      gap:10px;
      margin-bottom:14px;
    }}
    .turma-card__metrics div,
    .heatmap-card__meta-item {{
      padding:12px 14px;
      border-radius:18px;
      background:#f6f8fb;
      border:1px solid rgba(20,45,78,.06);
    }}
    .turma-card__metrics span,
    .heatmap-card__meta-item span {{
      display:block;
      font-size:11px;
      letter-spacing:.12em;
      text-transform:uppercase;
      color:var(--muted);
      margin-bottom:6px;
    }}
    .turma-card__metrics strong,
    .heatmap-card__meta-item strong {{
      color:var(--navy);
      font-size:17px;
    }}
    .turma-card__focus {{
      display:flex;
      flex-direction:column;
      gap:8px;
    }}
    .strength-stack {{
      display:grid;
      gap:12px;
      margin-top:12px;
    }}
    .strength-chip {{
      display:flex;
      justify-content:space-between;
      gap:12px;
      padding:14px 16px;
      border-radius:18px;
      background:linear-gradient(90deg, rgba(43,122,120,.12), rgba(214,170,61,.15));
      border:1px solid rgba(20,45,78,.08);
      font-weight:700;
    }}
    .heat-cell {{
      color:white;
      font-weight:700;
    }}
    .heatmap-grid {{
      display:grid;
      grid-template-columns:1fr;
      gap:16px;
    }}
    .heatmap-card {{
      padding:18px;
      border-radius:24px;
      background:#fffefb;
      border:1px solid rgba(20,45,78,.08);
      box-shadow:0 12px 28px rgba(18,33,56,.06);
    }}
    .heatmap-card__head {{
      display:grid;
      grid-template-columns:1.2fr .8fr;
      gap:16px;
      align-items:start;
      margin-bottom:14px;
    }}
    .heatmap-card__head h4 {{
      color:var(--navy);
      font-size:20px;
      margin-bottom:4px;
    }}
    .heatmap-card__head p {{
      color:var(--muted);
      line-height:1.45;
    }}
    .heatmap-card__meta {{
      display:grid;
      grid-template-columns:repeat(2, 1fr);
      gap:10px;
    }}
    .heatmap-card__grid {{
      display:grid;
      grid-template-columns:repeat(4, minmax(0, 1fr));
      gap:10px;
    }}
    .heat-tile {{
      display:flex;
      flex-direction:column;
      align-items:flex-start;
      gap:12px;
      min-height:90px;
      padding:12px;
      border:none;
      border-radius:18px;
      cursor:pointer;
      position:relative;
      text-align:left;
      font-family:inherit;
      box-shadow:inset 0 0 0 1px rgba(255,255,255,.18);
    }}
    .heat-tile strong {{
      font-size:20px;
      line-height:1;
    }}
    .heat-tile::after {{
      content:attr(data-tooltip);
      position:absolute;
      left:0;
      top:calc(100% + 10px);
      width:260px;
      padding:12px 14px;
      border-radius:16px;
      background:rgba(20,45,78,.96);
      color:white;
      font-size:12px;
      line-height:1.5;
      box-shadow:0 16px 30px rgba(18,33,56,.22);
      opacity:0;
      transform:translateY(4px);
      pointer-events:none;
      transition:opacity .18s ease, transform .18s ease;
      z-index:20;
      white-space:normal;
    }}
    .heat-tile:hover::after,
    .heat-tile:focus::after,
    .heat-tile:focus-visible::after {{
      opacity:1;
      transform:translateY(0);
    }}
    .heat-tile:focus-visible {{
      outline:2px solid var(--navy);
      outline-offset:2px;
    }}
    .actions-grid {{
      grid-template-columns:repeat(4,1fr);
    }}
    .action-card {{
      padding:20px;
      border-radius:24px;
      background:linear-gradient(180deg, rgba(20,45,78,.96), rgba(35,92,158,.92));
      color:white;
      box-shadow:var(--shadow);
    }}
    .action-card h4 {{
      margin-bottom:10px;
      font-size:20px;
    }}
    .action-card p {{
      color:rgba(255,255,255,.82);
      line-height:1.6;
    }}
    .footer {{
      margin-top:18px;
      padding:20px 24px;
      color:var(--muted);
      font-size:13px;
      line-height:1.6;
    }}
    @media (max-width: 1080px) {{
      .hero, .overview-grid, .priority-grid {{ grid-template-columns:1fr; }}
      .actions-grid {{ grid-template-columns:1fr 1fr; }}
      .school-profile__topline, .school-profile__skills {{ grid-template-columns:1fr; }}
      .heatmap-card__head {{ grid-template-columns:1fr; }}
      .heatmap-card__grid {{ grid-template-columns:repeat(4, minmax(0, 1fr)); }}
    }}
    @media (max-width: 760px) {{
      .page {{ width:min(100vw - 20px, 100%); margin:10px auto 28px; }}
      .hero {{ padding:24px; }}
      .hero h1 {{ font-size:34px; }}
      .hero-panel, .actions-grid, .school-meta, .level-legend, .turma-card__metrics, .heatmap-card__meta {{ grid-template-columns:1fr; }}
      .priority-row {{ grid-template-columns:1fr; }}
      .section {{ padding:20px; }}
      .nav {{ display:grid; grid-template-columns:1fr 1fr; }}
      .school-tablist {{ grid-template-columns:1fr; }}
      .turma-card-grid {{ grid-template-columns:1fr; }}
      .heatmap-card__grid {{ grid-template-columns:repeat(2, minmax(0, 1fr)); }}
      .heat-tile {{ min-height:84px; padding:10px; gap:10px; }}
      .heat-tile strong {{ font-size:18px; }}
      .skill-chip__label {{ white-space:normal; }}
      .skill-chip__sub {{ font-size:10px; }}
      .turma-card__badges {{ padding-left:0; }}
    }}
  </style>
</head>
<body>
  <main class="page">
    <section class="hero">
      <div>
        <p class="eyebrow">Secretaria Municipal de Educação de {html.escape(str(metadata["Município"]))}</p>
        <h1>AvaliaRJ 2025 · Painel Executivo do 2º Ano</h1>
        <p class="lead">
          Síntese visual para apresentação institucional do recorte de {html.escape(str(context["analysis_scope"]))}, com foco em participação, proficiência, distribuição de desempenho e habilidades prioritárias da rede municipal.
        </p>
        <div class="hero-tags">
          <span>{html.escape(str(metadata["Avaliação"]))}</span>
          <span>{html.escape(str(metadata["Ano Escolar"]))}</span>
          <span>{html.escape(str(metadata["Componente Curricular"]))}</span>
          <span>{html.escape(str(context["analysis_scope"]))}</span>
          <span>Gerado em {context["generated_at"].strftime("%d/%m/%Y %H:%M")}</span>
        </div>
      </div>
      <div class="hero-panel">
        <div class="kpi"><span>Participação</span><strong>{pt_percent(float(context["participation"]))}</strong><small>{pt_number(int(context["total_avaliados"]))} avaliados de {pt_number(int(context["total_previstos"]))} previstos</small></div>
        <div class="kpi"><span>Proficiência média</span><strong>{pt_number(float(context["network_prof"]), 2)}</strong><small>Resultado ponderado da rede</small></div>
        <div class="kpi"><span>Esperado+</span><strong>{pt_percent(float(context["expected"]))}</strong><small>Adequado + Avançado</small></div>
        <div class="kpi"><span>Crítico</span><strong>{pt_percent(float(context["critical"]))}</strong><small>Abaixo do básico + Básico</small></div>
      </div>
    </section>

    <nav class="nav">
      <a href="#rede">Rede</a>
      <a href="#escolas">Escolas</a>
      <a href="#habilidades">Habilidades</a>
      <a href="#acoes">Encaminhamentos</a>
    </nav>

    <section id="rede" class="section">
      <div class="section-head">
        <div>
          <h2>Panorama da Rede</h2>
          <p>Leitura rápida do desempenho municipal para tomada de decisão da Secretaria. {html.escape(str(context["analysis_scope_note"]))}</p>
        </div>
      </div>
      <div class="overview-grid">
        <div class="levels-panel">
          <p class="mini-label">Distribuição dos estudantes por faixa de desempenho</p>
          {''.join(level_blocks)}
        </div>
        <div class="story-panel">
          <p class="mini-label">Mensagens centrais</p>
          <div class="story-list">
            <div class="story-item"><strong>Destaque da rede</strong><span>Melhor escola: {html.escape(str(context["best_school"]["Escola"]))} com proficiência {pt_number(float(context["best_school"]["Proficiência Média"]))}.</span></div>
            <div class="story-item"><strong>Ponto de atenção</strong><span>Maior vulnerabilidade em {html.escape(str(context["worst_turma"]["Turma"]))}, da escola {html.escape(str(context["worst_turma"]["Escola"]))}, com proficiência {pt_number(float(context["worst_turma"]["Proficiência Média"]))}.</span></div>
            <div class="story-item"><strong>Potencial de escala</strong><span>Turma referência: {html.escape(str(context["best_turma"]["Turma"]))} / {html.escape(str(context["best_turma"]["Escola"]))}, com proficiência {pt_number(float(context["best_turma"]["Proficiência Média"]))}.</span></div>
          </div>
        </div>
      </div>
    </section>

    <section id="escolas" class="section">
      <div class="section-head">
        <div>
          <h2>Aprendizado por Escola</h2>
          <p>Leitura inspirada no QEdu: aprendizado adequado, faixas de interpretação, distribuição por níveis e habilidades prioritárias de cada escola.</p>
        </div>
      </div>
      <div class="school-tabs">
        <div class="school-tablist" role="tablist" aria-label="Escolas da rede">
          {''.join(school_tabs)}
        </div>
        <div class="school-panels">
          {''.join(school_panels)}
        </div>
      </div>
    </section>

    <section id="habilidades" class="section">
      <div class="section-head">
        <div>
          <h2>Habilidades da Rede</h2>
          <p>Recorte das habilidades prioritárias da rede, com heatmap geral municipal. As leituras por turma ficam disponíveis dentro da escola selecionada.</p>
        </div>
      </div>
      <div class="priority-grid">
        <div class="priority-panel">
          <p class="mini-label">Menor domínio na rede</p>
          {''.join(priority_rows)}
        </div>
        <div class="priority-panel">
          <p class="mini-label">Habilidades consolidadas</p>
          <div class="strength-stack">{''.join(strength_rows)}</div>
        </div>
      </div>
      <div class="heat-panel" style="margin-top:18px;">
        <p class="mini-label">Heatmap geral da rede</p>
        <div class="heatmap-grid">
          {network_heat_card}
        </div>
      </div>
    </section>

    <section id="acoes" class="section">
      <div class="section-head">
        <div>
          <h2>Encaminhamentos Sugeridos</h2>
          <p>Proposta de leitura executiva para a fala da Secretaria e para o plano de ação imediato.</p>
        </div>
      </div>
      <div class="actions-grid">
        {''.join(recommendations_html)}
      </div>
    </section>

  </main>
  <script>
    (() => {{
      const tabs = Array.from(document.querySelectorAll("[data-school-tab]"));
      const panels = Array.from(document.querySelectorAll("[data-school-panel]"));
      if (!tabs.length || !panels.length) {{
        return;
      }}

      const activateSchoolTab = (panelId) => {{
        tabs.forEach((tab, index) => {{
          const isActive = tab.dataset.schoolTab === panelId;
          tab.classList.toggle("active", isActive);
          tab.setAttribute("aria-selected", isActive ? "true" : "false");
          tab.tabIndex = isActive ? 0 : -1;
          if (isActive) {{
            tabs[index].focus({{ preventScroll: true }});
          }}
        }});

        panels.forEach((panel) => {{
          const isActive = panel.dataset.schoolPanel === panelId;
          panel.classList.toggle("active", isActive);
          if (isActive) {{
            panel.removeAttribute("hidden");
          }} else {{
            panel.setAttribute("hidden", "hidden");
          }}
        }});
      }};

      tabs.forEach((tab, index) => {{
        tab.tabIndex = index === 0 ? 0 : -1;
        tab.addEventListener("click", () => activateSchoolTab(tab.dataset.schoolTab));
        tab.addEventListener("keydown", (event) => {{
          if (!["ArrowRight", "ArrowLeft", "ArrowDown", "ArrowUp", "Home", "End"].includes(event.key)) {{
            return;
          }}
          event.preventDefault();
          let targetIndex = index;
          if (event.key === "ArrowRight" || event.key === "ArrowDown") {{
            targetIndex = (index + 1) % tabs.length;
          }}
          if (event.key === "ArrowLeft" || event.key === "ArrowUp") {{
            targetIndex = (index - 1 + tabs.length) % tabs.length;
          }}
          if (event.key === "Home") {{
            targetIndex = 0;
          }}
          if (event.key === "End") {{
            targetIndex = tabs.length - 1;
          }}
          activateSchoolTab(tabs[targetIndex].dataset.schoolTab);
        }});
      }});
    }})();
  </script>
</body>
</html>"""

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(html_text, encoding="utf-8")
    return output_path


def add_textbox(slide, left, top, width, height, text, font_size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    paragraph.alignment = align
    return box


def add_card(slide, left, top, width, height, fill_color, line_color=None, radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_metric_card(slide, left, top, width, height, title, value, subtitle, fill_color, value_color=INK):
    add_card(slide, left, top, width, height, fill_color)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.22), title, 10, NAVY, True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.40), width - Inches(0.36), Inches(0.34), value, 22, value_color, True)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.82), width - Inches(0.36), Inches(0.34), subtitle, 10, INK)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.2), Inches(0.22), text, 9, RGBColor(93, 105, 122))


def draw_school_learning_panel(slide, profile: dict[str, object], left, top, width, height):
    add_card(slide, left, top, width, height, RGBColor(255, 255, 255))
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.3), str(profile["school"]), 12, NAVY, True)
    band = profile["band"]
    tone_hex = str(band["tone"]).lstrip("#")
    tone = RGBColor(int(tone_hex[0:2], 16), int(tone_hex[2:4], 16), int(tone_hex[4:6], 16))
    add_card(slide, left + Inches(0.18), top + Inches(0.55), Inches(1.55), Inches(0.35), tone)
    add_textbox(slide, left + Inches(0.27), top + Inches(0.62), Inches(1.35), Inches(0.16), str(band["label"]), 9, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.0), Inches(1.5), Inches(0.22), "Aprendizado adequado", 9, RGBColor(97, 112, 132), True)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.22), Inches(1.5), Inches(0.45), pt_percent(float(profile["adequate"])), 22, NAVY, True)
    add_textbox(slide, left + Inches(1.95), top + Inches(1.05), Inches(1.2), Inches(0.22), "Proficiência", 9, RGBColor(97, 112, 132), True)
    add_textbox(slide, left + Inches(1.95), top + Inches(1.28), Inches(1.1), Inches(0.3), pt_number(float(profile["prof"])), 16, BLUE, True)
    add_textbox(slide, left + Inches(3.0), top + Inches(1.05), Inches(1.0), Inches(0.22), "Participação", 9, RGBColor(97, 112, 132), True)
    add_textbox(slide, left + Inches(3.0), top + Inches(1.28), Inches(1.0), Inches(0.3), pt_percent(float(profile["participation"])), 16, INK, True)

    scale_left = left + Inches(0.2)
    scale_top = top + Inches(1.78)
    scale_width = width - Inches(0.4)
    segments = [
        (0.25, RGBColor(201, 92, 84)),
        (0.25, RGBColor(216, 137, 77)),
        (0.20, RGBColor(210, 178, 76)),
        (0.30, RGBColor(47, 133, 90)),
    ]
    running_left = scale_left
    for fraction, color in segments:
        seg_width = scale_width * fraction
        add_card(slide, running_left, scale_top, seg_width, Inches(0.12), color, color, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        running_left += seg_width
    marker_offset = scale_width * (float(profile["adequate"]) / 100.0)
    add_card(slide, scale_left + marker_offset, scale_top - Inches(0.05), Inches(0.03), Inches(0.24), NAVY, NAVY, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    for label, xpos in zip(["0", "25", "50", "70", "100"], [0.0, 0.25, 0.5, 0.7, 1.0]):
        add_textbox(slide, scale_left + scale_width * xpos - Inches(0.08), scale_top + Inches(0.12), Inches(0.18), Inches(0.16), label, 7, RGBColor(107, 114, 128), False, PP_ALIGN.CENTER)

    add_textbox(slide, left + Inches(0.18), top + Inches(2.18), Inches(2.2), Inches(0.2), "Distribuição do aprendizado", 9, RGBColor(97, 112, 132), True)
    stack_left = left + Inches(0.18)
    stack_top = top + Inches(2.46)
    stack_width = width - Inches(0.36)
    colors = {
        "Abaixo do básico": RGBColor(201, 92, 84),
        "Básico": RGBColor(216, 137, 77),
        "Adequado": RGBColor(210, 178, 76),
        "Avançado": RGBColor(47, 133, 90),
    }
    run = stack_left
    for label, value in profile["distribution"]:
        seg_width = stack_width * (float(value) / 100.0)
        add_card(slide, run, stack_top, seg_width, Inches(0.16), colors[label], colors[label], MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        run += seg_width
    legend_top = top + Inches(2.72)
    legend_positions = [(0.18, "Abaixo do básico"), (1.45, "Básico"), (2.35, "Adequado"), (3.35, "Avançado")]
    dist_map = {label: value for label, value in profile["distribution"]}
    for xoffset, label in legend_positions:
        add_card(slide, left + Inches(xoffset), legend_top + Inches(0.03), Inches(0.08), Inches(0.08), colors[label], colors[label], MSO_AUTO_SHAPE_TYPE.OVAL)
        add_textbox(slide, left + Inches(xoffset + 0.11), legend_top, Inches(0.95), Inches(0.14), label, 7.5, INK)
        add_textbox(slide, left + Inches(xoffset + 0.11), legend_top + Inches(0.12), Inches(0.7), Inches(0.14), pt_percent(float(dist_map[label])), 7.5, NAVY, True)

    add_textbox(slide, left + Inches(0.18), top + Inches(3.02), Inches(1.7), Inches(0.18), "Prioridades", 9, RGBColor(97, 112, 132), True)
    add_textbox(slide, left + Inches(2.1), top + Inches(3.02), Inches(1.7), Inches(0.18), "Forças", 9, RGBColor(97, 112, 132), True)
    y = top + Inches(3.22)
    for label, value in profile["weak_skills"][:2]:
        add_textbox(slide, left + Inches(0.18), y, Inches(1.75), Inches(0.16), f"• {label} {pt_percent(float(value))}", 8.5, RED)
        y += Inches(0.22)
    y = top + Inches(3.22)
    for label, value in profile["strong_skills"][:2]:
        add_textbox(slide, left + Inches(2.1), y, Inches(1.75), Inches(0.16), f"• {label} {pt_percent(float(value))}", 8.5, GREEN)
        y += Inches(0.22)


def generate_powerpoint(context: dict[str, object], output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    metadata = context["metadata"]
    school_report = context["school_report"]
    turma_report = context["turma_report"]
    skill_summary = context["skill_summary"]
    school_profiles = context["school_profiles"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_card(slide, Inches(8.9), Inches(-0.35), Inches(5.5), Inches(8.5), BLUE)
    add_card(slide, Inches(9.8), Inches(4.7), Inches(3.0), Inches(2.2), GOLD)
    add_textbox(slide, Inches(0.65), Inches(0.72), Inches(7.4), Inches(0.4), "Secretaria Municipal de Educação de Mendes", 14, RGBColor(221, 229, 240), True)
    add_textbox(slide, Inches(0.65), Inches(1.45), Inches(7.2), Inches(1.2), "AvaliaRJ 2025\nPainel Executivo do 2º Ano", 24, RGBColor(255, 255, 255), True)
    add_textbox(
        slide,
        Inches(0.65),
        Inches(3.0),
        Inches(6.8),
        Inches(1.2),
        (
            f"Leitura executiva do recorte {metadata['Componente Curricular']} · "
            f"{metadata['Ano Escolar']} · {metadata['Município']}."
        ),
        16,
        RGBColor(230, 236, 244),
    )
    add_metric_card(slide, Inches(0.65), Inches(4.75), Inches(2.2), Inches(1.3), "Participação", pt_percent(float(context["participation"])), "rede avaliada", RGBColor(241, 245, 249))
    add_metric_card(slide, Inches(3.0), Inches(4.75), Inches(2.2), Inches(1.3), "Proficiência", pt_number(float(context["network_prof"]), 2), "média ponderada", RGBColor(241, 245, 249))
    add_metric_card(slide, Inches(5.35), Inches(4.75), Inches(2.2), Inches(1.3), "Esperado+", pt_percent(float(context["expected"])), "adequado + avançado", RGBColor(241, 245, 249))
    add_footer(slide, f"Gerado em {context['generated_at'].strftime('%d/%m/%Y %H:%M')} a partir de arquivos agregados do AvaliaRJ.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SAND
    add_textbox(slide, Inches(0.55), Inches(0.4), Inches(7.5), Inches(0.45), "Panorama da Rede", 24, NAVY, True)
    add_textbox(slide, Inches(0.55), Inches(0.85), Inches(8.0), Inches(0.35), "Síntese para abertura da apresentação institucional.", 11, RGBColor(89, 98, 114))
    add_metric_card(slide, Inches(0.55), Inches(1.35), Inches(2.0), Inches(1.15), "Escolas", pt_number(school_report['Escola'].nunique()), "unidades com resultado", RGBColor(255, 255, 255))
    add_metric_card(slide, Inches(2.75), Inches(1.35), Inches(2.0), Inches(1.15), "Turmas", pt_number(turma_report['Turma'].nunique()), "turmas avaliadas", RGBColor(255, 255, 255))
    add_metric_card(slide, Inches(4.95), Inches(1.35), Inches(2.0), Inches(1.15), "Crítico", pt_percent(float(context["critical"])), "abaixo do básico + básico", RGBColor(255, 255, 255))
    add_metric_card(slide, Inches(7.15), Inches(1.35), Inches(2.0), Inches(1.15), "Esperado+", pt_percent(float(context["expected"])), "adequado + avançado", RGBColor(255, 255, 255))

    add_card(slide, Inches(0.55), Inches(2.8), Inches(4.2), Inches(3.5), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(3.6), Inches(0.3), "Distribuição de desempenho", 14, NAVY, True)
    level_top = Inches(3.45)
    for label, value in context["levels"].items():
        add_textbox(slide, Inches(0.85), level_top, Inches(1.5), Inches(0.25), label, 11, INK, True)
        add_textbox(slide, Inches(3.75), level_top, Inches(0.7), Inches(0.25), pt_percent(float(value)), 11, INK, True, PP_ALIGN.RIGHT)
        bar_bg = add_card(slide, Inches(1.9), level_top + Inches(0.03), Inches(1.7), Inches(0.12), RGBColor(226, 232, 240), RGBColor(226, 232, 240), MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        bar_fill_width = Inches(1.7 * float(value) / 100.0)
        add_card(slide, Inches(1.9), level_top + Inches(0.03), bar_fill_width, Inches(0.12), TEAL if label in {"Adequado", "Avançado"} else GOLD, TEAL if label in {"Adequado", "Avançado"} else GOLD, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        level_top += Inches(0.55)

    add_card(slide, Inches(5.1), Inches(2.8), Inches(3.35), Inches(1.45), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(5.35), Inches(3.0), Inches(2.8), Inches(0.25), "Melhor escola", 12, NAVY, True)
    add_textbox(slide, Inches(5.35), Inches(3.35), Inches(2.8), Inches(0.6), str(context["best_school"]["Escola"]), 18, BLUE, True)
    add_textbox(slide, Inches(5.35), Inches(3.88), Inches(2.8), Inches(0.25), f"Proficiência {pt_number(float(context['best_school']['Proficiência Média']))}", 11, INK)

    add_card(slide, Inches(5.1), Inches(4.45), Inches(3.35), Inches(1.45), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(5.35), Inches(4.65), Inches(2.8), Inches(0.25), "Turma de atenção imediata", 12, NAVY, True)
    add_textbox(slide, Inches(5.35), Inches(5.0), Inches(2.8), Inches(0.6), f"{context['worst_turma']['Turma']} · {context['worst_turma']['Escola']}", 16, RED, True)
    add_textbox(slide, Inches(5.35), Inches(5.5), Inches(2.8), Inches(0.25), f"Proficiência {pt_number(float(context['worst_turma']['Proficiência Média']))}", 11, INK)

    add_card(slide, Inches(8.7), Inches(2.8), Inches(4.0), Inches(3.1), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(8.95), Inches(3.0), Inches(3.4), Inches(0.25), "Mensagem-chave para a fala", 12, NAVY, True)
    messages = [
        f"A rede manteve participação de {pt_percent(float(context['participation']))}.",
        f"O resultado médio foi {pt_number(float(context['network_prof']), 2)} de proficiência.",
        f"As fragilidades concentram-se em {', '.join(item['Habilidade'] for item in context['priority_skills'])}.",
        f"Há referência positiva em {context['best_school']['Escola']} e na turma {context['best_turma']['Turma']}.",
    ]
    bullet_top = Inches(3.45)
    for message in messages:
        add_textbox(slide, Inches(9.05), bullet_top, Inches(3.2), Inches(0.45), "• " + message, 12, INK)
        bullet_top += Inches(0.62)
    add_footer(slide, "Panorama institucional para abertura da reunião com a Secretaria.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 248, 244)
    add_textbox(slide, Inches(0.55), Inches(0.4), Inches(8.2), Inches(0.45), "Aprendizado por Escola", 24, NAVY, True)
    add_textbox(slide, Inches(0.55), Inches(0.85), Inches(10.0), Inches(0.3), "Leitura inspirada no QEdu: aprendizado adequado, faixas de interpretação, distribuição por níveis e habilidades foco.", 11, RGBColor(89, 98, 114))
    panel_lefts = [0.55, 4.5, 8.45]
    for profile, left in zip(school_profiles, panel_lefts):
        draw_school_learning_panel(slide, profile, Inches(left), Inches(1.35), Inches(3.8), Inches(5.35))
    add_footer(slide, "Painéis por escola com foco em aprendizado adequado e leitura pedagógica do resultado.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SAND
    add_textbox(slide, Inches(0.55), Inches(0.4), Inches(7.0), Inches(0.45), "Resultados por Turma", 24, NAVY, True)
    add_textbox(slide, Inches(0.55), Inches(0.85), Inches(8.5), Inches(0.3), "Recorte útil para acompanhamento pedagógico e apoio focalizado.", 11, RGBColor(89, 98, 114))
    add_card(slide, Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.9), RGBColor(255, 255, 255))
    headers = ["Rank", "Turma", "Escola", "Prof.", "Part.", "Esperado+", "Foco"]
    col_left = [0.75, 1.45, 3.0, 6.5, 7.45, 8.45, 9.7]
    widths = [0.5, 1.4, 3.3, 0.8, 0.8, 1.0, 2.6]
    for idx, header in enumerate(headers):
        add_textbox(slide, Inches(col_left[idx]), Inches(1.6), Inches(widths[idx]), Inches(0.25), header, 10, NAVY, True)

    row_top = Inches(1.95)
    for _, row in turma_report.iterrows():
        band_color = RGBColor(248, 250, 252) if int(row["Ranking"]) % 2 else RGBColor(241, 245, 249)
        add_card(slide, Inches(0.7), row_top - Inches(0.05), Inches(11.8), Inches(0.58), band_color, band_color, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
        values = [
            str(int(row["Ranking"])),
            str(row["Turma"]),
            str(row["Escola"]),
            pt_number(float(row["Proficiência Média"])),
            pt_percent(float(row["Participacao %"])),
            pt_percent(float(row["Esperado+ (%)"])),
            str(row["Habilidade mais frágil"]),
        ]
        for idx, value in enumerate(values):
            add_textbox(slide, Inches(col_left[idx]), row_top, Inches(widths[idx]), Inches(0.22), value, 9 if idx != 2 else 8.5, INK, idx == 0)
        row_top += Inches(0.62)

    add_card(slide, Inches(9.3), Inches(6.45), Inches(3.45), Inches(0.55), GOLD, GOLD)
    add_textbox(slide, Inches(9.45), Inches(6.58), Inches(3.0), Inches(0.2), f"Destaque: {context['best_turma']['Turma']} · {pt_number(float(context['best_turma']['Proficiência Média']))}", 10, NAVY, True)
    add_footer(slide, "Todas as turmas do recorte foram incluídas para leitura de apoio à coordenação pedagógica.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 248, 244)
    add_textbox(slide, Inches(0.55), Inches(0.4), Inches(7.5), Inches(0.45), "Habilidades Prioritárias", 24, NAVY, True)
    add_textbox(slide, Inches(0.55), Inches(0.85), Inches(9.0), Inches(0.3), "As habilidades Hxx abaixo concentram a principal agenda de recomposição da rede.", 11, RGBColor(89, 98, 114))

    skill_chart_data = CategoryChartData()
    skill_chart_data.categories = list(skill_summary["Habilidade"])
    skill_chart_data.add_series("Média da rede (%)", list(skill_summary["Media da rede (%)"]))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(1.45), Inches(7.0), Inches(4.9), skill_chart_data).chart
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = TEAL
    chart.has_legend = False
    chart.category_axis.reverse_order = True
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Domínio médio por habilidade"

    add_card(slide, Inches(7.9), Inches(1.45), Inches(4.85), Inches(2.1), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(8.15), Inches(1.68), Inches(4.1), Inches(0.25), "Prioridades imediatas", 12, NAVY, True)
    box_top = Inches(2.1)
    for item in context["priority_skills"]:
        add_textbox(slide, Inches(8.15), box_top, Inches(2.2), Inches(0.25), str(item["Habilidade"]), 12, RED, True)
        add_textbox(slide, Inches(10.6), box_top, Inches(1.2), Inches(0.25), pt_percent(float(item["Media da rede (%)"])), 12, INK, True, PP_ALIGN.RIGHT)
        box_top += Inches(0.42)

    add_card(slide, Inches(7.9), Inches(3.8), Inches(4.85), Inches(2.05), RGBColor(255, 255, 255))
    add_textbox(slide, Inches(8.15), Inches(4.0), Inches(4.1), Inches(0.25), "Forças a preservar", 12, NAVY, True)
    box_top = Inches(4.42)
    for item in context["strengths"]:
        add_textbox(slide, Inches(8.15), box_top, Inches(2.2), Inches(0.25), str(item["Habilidade"]), 12, GREEN, True)
        add_textbox(slide, Inches(10.6), box_top, Inches(1.2), Inches(0.25), pt_percent(float(item["Media da rede (%)"])), 12, INK, True, PP_ALIGN.RIGHT)
        box_top += Inches(0.42)
    add_footer(slide, "Leitura de habilidades com legenda documental de captura 2025 quando disponível.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, Inches(0.55), Inches(0.5), Inches(8.5), Inches(0.45), "Encaminhamentos para a Secretaria", 24, RGBColor(255, 255, 255), True)
    add_textbox(slide, Inches(0.55), Inches(0.95), Inches(9.0), Inches(0.3), "Síntese final orientada à ação pedagógica e de gestão.", 11, RGBColor(214, 224, 236))
    card_positions = [(0.65, 1.55), (6.75, 1.55), (0.65, 4.1), (6.75, 4.1)]
    for item, (left, top) in zip(context["recommendations"], card_positions):
        add_card(slide, Inches(left), Inches(top), Inches(5.2), Inches(2.1), RGBColor(246, 249, 252))
        add_textbox(slide, Inches(left + 0.22), Inches(top + 0.18), Inches(4.7), Inches(0.35), item["title"], 16, NAVY, True)
        add_textbox(slide, Inches(left + 0.22), Inches(top + 0.62), Inches(4.7), Inches(1.2), item["body"], 12, INK)
    add_footer(slide, "Material de apoio à apresentação institucional. Uso pontual, sem integração ao frontend oficial do produto.")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(output_path)
        return output_path
    except PermissionError:
        alt_path = output_path.with_name(
            f"{output_path.stem}_{datetime.now():%Y%m%d-%H%M%S}{output_path.suffix}"
        )
        prs.save(alt_path)
        return alt_path


def main() -> None:
    args = parse_args()
    context = build_context(args)
    html_path = render_dashboard_html(context, Path(args.saida_html))
    ppt_path = generate_powerpoint(context, Path(args.saida_pptx))
    print(f"Dashboard HTML gerado em: {html_path.resolve()}")
    print(f"Apresentacao PowerPoint gerada em: {ppt_path.resolve()}")


if __name__ == "__main__":
    main()
